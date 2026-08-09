#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""check_deck_freshness.py — has the data moved past the committed MCOM deck?

WHY THIS EXISTS
---------------
build_mcom_macro_pptx.py reads every figure out of platform/data at build time, which is the right
design: the deck cannot drift from the platform by being edited, because it is never edited. But it
CAN drift by standing still. The data refreshes on a dozen crons; the deck only refreshes when
somebody remembers to run the script. Nothing in .github/workflows references docs/decks at all, so
"somebody remembers" was the entire mechanism.

That gap is not hypothetical. The committed deck was built on 2026-08-05; by 2026-08-09 a rebuild
differed in 62 lines — USD/THB 33.34 -> 33.05, rivers running high 132/795 -> 161/796, cassava
+59.3% -> +61.3%, palm +27.3% -> +25.5%, and so on. Every one of those is a number somebody could
have read aloud in a meeting, from a deck that looked authoritative and was simply out of date.

WHAT IT DOES
------------
Rebuilds the deck into a temp directory and compares the rendered TEXT of every slide against the
committed .pptx. Text, not bytes: a .pptx is a zip, so two byte-different files can say exactly the
same thing (timestamps, zip ordering, part ids), and a byte comparison would cry wolf on every run.
What matters to a reader is what the slide says.

Reports the specific figures that moved, slide by slide, so the diff is reviewable at a glance
rather than being "the deck changed, rebuild it".

DETERMINISTIC + NETWORK-FREE — it builds from committed data and reads a committed file.

Exit codes match the gate idiom used by check_commodity_board.py / check_feed_liveness.py:
  0  the committed deck still says what the data says — nothing to do
  2  the data has moved past the deck; the committed deck is showing stale figures
  3  cannot judge — python-pptx missing, the build failed, or the committed deck is absent
"""
from __future__ import annotations

import argparse
import difflib
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
BUILDER = HERE / "build_mcom_macro_pptx.py"

# The builder hardcodes its output name (the date in it is the deck's edition, not the data
# vintage), so the committed artefact is addressed by the same literal rather than by globbing —
# a glob would silently pick up an older edition left in the folder.
DECK_NAME = "mcom-2026-08-05-macro.pptx"


def slide_text(path: Path):
    """Every non-empty text run, grouped by slide.

    Runs rather than whole shapes: a figure and its label usually live in separate runs, so run
    granularity is what makes the diff read as "this number changed" instead of "this paragraph
    changed".
    """
    from pptx import Presentation  # imported here so a missing dep exits 3, not a traceback

    out = []
    for slide in Presentation(str(path)).slides:
        runs = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = (run.text or "").strip()
                    if t:
                        runs.append(t)
        out.append(runs)
    return out


def build_fresh(dest: Path, lang: str) -> Path:
    """Run the real builder, so this check can never disagree with what a rebuild would produce."""
    cmd = [sys.executable, str(BUILDER), "--out", str(dest)]
    if lang != "en":
        cmd += ["--lang", lang]
    proc = subprocess.run(cmd, cwd=str(HERE), capture_output=True, text=True)
    if proc.returncode != 0:
        sys.stderr.write("check_deck_freshness.py: the deck build failed (rc=%d)\n" % proc.returncode)
        sys.stderr.write((proc.stderr or proc.stdout or "")[-2000:] + "\n")
        return None
    suffix = "" if lang == "en" else "-TH"
    built = dest / DECK_NAME.replace(".pptx", "%s.pptx" % suffix)
    return built if built.exists() else None


def report(committed_slides, fresh_slides):
    """Per-slide line diff. Returns the list of human-readable changes."""
    changes = []
    n = max(len(committed_slides), len(fresh_slides))
    for i in range(n):
        old = committed_slides[i] if i < len(committed_slides) else []
        new = fresh_slides[i] if i < len(fresh_slides) else []
        if old == new:
            continue
        if not old:
            changes.append("slide %d: NEW slide in the rebuild (%d lines)" % (i + 1, len(new)))
            continue
        if not new:
            changes.append("slide %d: slide disappears in the rebuild" % (i + 1))
            continue
        sm = difflib.SequenceMatcher(a=old, b=new, autojunk=False)
        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            if tag == "equal":
                continue
            if tag == "replace" and (i2 - i1) == (j2 - j1):
                for k in range(i2 - i1):          # a clean 1:1 swap is the "a figure moved" case
                    changes.append("slide %d: %r -> %r" % (i + 1, old[i1 + k], new[j1 + k]))
            elif tag == "delete":
                for t in old[i1:i2]:
                    changes.append("slide %d: dropped %r" % (i + 1, t))
            elif tag == "insert":
                for t in new[j1:j2]:
                    changes.append("slide %d: added %r" % (i + 1, t))
            else:
                changes.append("slide %d: %r -> %r"
                               % (i + 1, " | ".join(old[i1:i2]), " | ".join(new[j1:j2])))
    return changes


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--lang", default="en", help="en (default) or th")
    ap.add_argument("--write", action="store_true",
                    help="when the data has moved, overwrite the committed deck with the rebuild "
                         "(what the workflow does before opening its PR)")
    ap.add_argument("--check", action="store_true", help="accepted for symmetry with the other guards")
    a = ap.parse_args(argv)

    suffix = "" if a.lang == "en" else "-TH"
    committed = HERE / DECK_NAME.replace(".pptx", "%s.pptx" % suffix)
    if not committed.exists():
        sys.stderr.write("check_deck_freshness.py: committed deck %s is absent — nothing to compare "
                         "against.\n" % committed.name)
        return 3
    try:
        import pptx  # noqa: F401
    except ImportError:
        sys.stderr.write("check_deck_freshness.py: python-pptx is not installed — cannot read the "
                         "deck. pip install python-pptx\n")
        return 3

    with tempfile.TemporaryDirectory() as td:
        fresh = build_fresh(Path(td), a.lang)
        if fresh is None:
            return 3
        try:
            committed_slides = slide_text(committed)
            fresh_slides = slide_text(fresh)
        except Exception as exc:                      # a corrupt pptx is "cannot judge", not "stale"
            sys.stderr.write("check_deck_freshness.py: could not read a deck: %s\n" % exc)
            return 3

        changes = report(committed_slides, fresh_slides)
        print("check_deck_freshness.py: %s — %d slides committed, %d rebuilt"
              % (committed.name, len(committed_slides), len(fresh_slides)))

        if not changes:
            print("check_deck_freshness.py: OK — the committed deck still says what the data says.")
            return 0

        print("\nThe data has moved past the committed deck. %d line(s) differ:\n" % len(changes))
        for c in changes:
            print("  " + c)

        if a.write:
            # Copied as bytes: the rebuild IS the new deck, and re-running the builder straight
            # over the committed path would build it a second time for no reason.
            committed.write_bytes(fresh.read_bytes())
            print("\ncheck_deck_freshness.py: rewrote %s from the rebuild." % committed.name)
        else:
            print("\nRebuild with:  python docs/decks/build_mcom_macro_pptx.py"
                  + ("" if a.lang == "en" else "  --lang th"))
        return 2


if __name__ == "__main__":
    sys.exit(main())
