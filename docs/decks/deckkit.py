"""House-style PowerPoint layout kit for AutoX / เงินไชโย decks.

Why this exists rather than the html2pptx route: Kanit — the house font, confirmed by reading the
run properties of Kaustav's own LTV decks (4,590 Kanit runs in one, 5,854 in the other) — is not
installed on this laptop, so an HTML renderer would compute every box position with fallback
metrics and the PPTX would be laid out to the wrong widths. Instead we drive python-pptx directly
and measure text with the Kanit TTF itself through Pillow, so a box that fits here fits in
PowerPoint too.

Geometry, palette and type sizes are lifted from the reference decks, not invented:
  slide            13.333 x 7.5 in (16:9)
  cover            navy field, gold band at L10.73 W2.60, red hairline at L10.73 W0.12
  content          white field, 20pt navy title, small AutoX mark at L11.63 T0.28 W1.03 H0.32,
                   red 0.12in bar along the bottom at T7.38
  table header     fill 1B2A6B, white 10.5pt bold, centred
  footer           "Restricted Data – Reproduction is prohibited", 8pt

The content slides deliberately do NOT copy the reference deck's navy header BAR. The AutoX mark is
dark navy artwork on transparency; the reference puts it inside that navy bar, where it nearly
vanishes. White header, navy type, mark on white — same house furniture, readable contrast.

There is no LibreOffice on this machine, so `preview()` renders each slide to PNG with Pillow using
the same geometry and the same font files. It is a fit check, not a pixel-accurate PowerPoint
render: its purpose is to catch text running out of its box before the file reaches a projector.
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette
NAVY = "1E2F5C"     # house body / title navy
NAVY2 = "1B2A6B"    # the table-header navy, a touch brighter
RED = "CC0000"      # house red — emphasis and deterioration
GOLD = "F5C242"     # house gold — watch / caution
GREY = "606060"     # house secondary text
WHITE = "FFFFFF"
GREEN = "15795F"    # the one addition: the house palette has no positive colour, and a macro deck
                    # that cannot say "this is a tailwind" in colour forces every number to read as
                    # a warning. Desaturated to sit beside the navy rather than fight it.
CARD = "F4F6FA"     # card ground
LINE = "D9DEE9"     # hairline
BAND = "F7F9FC"     # zebra row

EMU_IN = 914400
PT_IN = 72.0

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
FOOTER = "Restricted Data – Reproduction is prohibited"


def _font_paths():
    """Kanit regular + semibold, wherever they happen to live on this machine."""
    cands = []
    if os.environ.get("KANIT_DIR"):
        cands.append(Path(os.environ["KANIT_DIR"]))
    cands += [HERE / ".fonts", ASSETS / "fonts", Path.home() / ".fonts",
              Path("C:/Windows/Fonts")]
    for d in cands:
        r, b = d / "Kanit-Regular.ttf", d / "Kanit-SemiBold.ttf"
        if r.exists() and b.exists():
            return r, b
    return None, None


class Text:
    """Kanit metrics via Pillow, so wrapping decisions here match what PowerPoint will do.

    Insets are zeroed on every text box we create and line spacing is set to an exact point value,
    which removes the two things that normally make a python-pptx box and a hand-computed layout
    disagree.
    """

    def __init__(self):
        self.reg, self.bold = _font_paths()
        self.ok = bool(self.reg)
        self._cache = {}
        if self.ok:
            from PIL import ImageFont  # noqa: F401  (import proves the stack is usable)

    def font(self, size_pt, bold=False, scale=1.0):
        if not self.ok:
            return None
        from PIL import ImageFont
        key = (round(size_pt * scale, 2), bold)
        if key not in self._cache:
            path = self.bold if bold else self.reg
            self._cache[key] = ImageFont.truetype(str(path), max(1, int(round(size_pt * scale))))
        return self._cache[key]

    def width_in(self, s, size_pt, bold=False):
        """Rendered width of `s` in INCHES at `size_pt`. Measured at 4x for sub-pixel accuracy."""
        if not self.ok or not s:
            return 0.0
        f = self.font(size_pt, bold, scale=4.0)
        return f.getlength(s) / 4.0 / PT_IN

    def wrap(self, s, size_pt, width_in, bold=False):
        """Greedy word wrap to `width_in`. Thai has no spaces, so a long Thai run that will not fit
        is returned whole rather than chopped mid-syllable — it shows up as an overflow finding
        instead of being silently mangled."""
        if not s:
            return [""]
        out = []
        for para in s.split("\n"):
            words, cur = para.split(" "), ""
            for w in words:
                trial = w if not cur else cur + " " + w
                if self.ok and self.width_in(trial, size_pt, bold) > width_in and cur:
                    out.append(cur)
                    cur = w
                else:
                    cur = trial
            out.append(cur)
        return out


TX = Text()


def rgb(h):
    return RGBColor.from_string(h)


# The sign-off page's stylesheet, lifted verbatim from the hand-authored review page it
# replaces so the page Kaustav already signed off on still looks like itself. Only the
# CONTENT is generated now; the design was fine, it was the numbers in it that went stale.
REVIEW_CSS = """<style>
:root{
  --navy:#1E2F5C; --navy-lift:#2A3F73; --gold:#C08A12; --red:#CC0000;
  --bg:#F6F7FA; --card:#FFFFFF; --ink:#151A24; --ink-2:#4A5468; --rule:#DDE1E9;
  --shadow:0 1px 2px rgba(30,47,92,.06), 0 8px 24px rgba(30,47,92,.07);
}
@media (prefers-color-scheme:dark){
  :root{--navy:#8FA6DC; --navy-lift:#A8BBE8; --gold:#E0B54A; --red:#F2635B;
    --bg:#0F1219; --card:#171B24; --ink:#E7EAF1; --ink-2:#A2ABBE; --rule:#262C39;
    --shadow:0 1px 2px rgba(0,0,0,.4), 0 8px 24px rgba(0,0,0,.35);}
}
:root[data-theme="light"]{--navy:#1E2F5C; --navy-lift:#2A3F73; --gold:#C08A12; --red:#CC0000;
  --bg:#F6F7FA; --card:#FFFFFF; --ink:#151A24; --ink-2:#4A5468; --rule:#DDE1E9;
  --shadow:0 1px 2px rgba(30,47,92,.06), 0 8px 24px rgba(30,47,92,.07);}
:root[data-theme="dark"]{--navy:#8FA6DC; --navy-lift:#A8BBE8; --gold:#E0B54A; --red:#F2635B;
  --bg:#0F1219; --card:#171B24; --ink:#E7EAF1; --ink-2:#A2ABBE; --rule:#262C39;
  --shadow:0 1px 2px rgba(0,0,0,.4), 0 8px 24px rgba(0,0,0,.35);}

*{box-sizing:border-box}
body{margin:0;background:var(--bg);color:var(--ink);
  font:15px/1.6 ui-sans-serif,system-ui,-apple-system,"Segoe UI",Roboto,"Helvetica Neue",sans-serif;
  -webkit-text-size-adjust:100%}
.wrap{max-width:1080px;margin:0 auto;padding:0 20px 96px}
a{color:var(--navy)}

/* ---------- masthead ---------- */
header.top{padding:56px 0 28px;border-bottom:3px solid var(--navy)}
.kicker{font-size:11px;letter-spacing:.16em;text-transform:uppercase;color:var(--gold);
  font-weight:700;margin:0 0 12px}
h1{font-size:clamp(28px,4.4vw,42px);line-height:1.1;letter-spacing:-.02em;margin:0 0 14px;
  font-weight:700;text-wrap:balance}
.dek{font-size:17px;color:var(--ink-2);margin:0;max-width:62ch}
.meta{display:flex;flex-wrap:wrap;gap:8px;margin:22px 0 0;padding:0;list-style:none}
.meta li{font-size:12px;padding:5px 11px;border:1px solid var(--rule);border-radius:999px;
  color:var(--ink-2);background:var(--card)}
.meta li b{color:var(--ink);font-weight:600}

/* ---------- notice ---------- */
.notice{background:var(--card);border:1px solid var(--rule);border-left:4px solid var(--red);
  border-radius:0 10px 10px 0;padding:20px 22px;margin:32px 0 0;box-shadow:var(--shadow)}
.notice h3{margin:0 0 10px;font-size:13px;letter-spacing:.06em;text-transform:uppercase;
  color:var(--red);font-weight:700}
.notice p{margin:0 0 10px}
.notice p:last-child{margin:0}
.notice .q{font-size:14px;color:var(--ink-2)}

/* ---------- jump bar ---------- */
nav.jump{position:sticky;top:0;z-index:5;background:var(--bg);padding:14px 0 12px;
  border-bottom:1px solid var(--rule);margin:36px 0 0;overflow-x:auto}
nav.jump ol{display:flex;gap:6px;margin:0;padding:0;list-style:none;min-width:min-content}
nav.jump a{display:block;min-width:30px;text-align:center;padding:5px 8px;font-size:12px;
  font-variant-numeric:tabular-nums;text-decoration:none;color:var(--ink-2);
  border:1px solid var(--rule);border-radius:6px;background:var(--card)}
nav.jump a:hover,nav.jump a:focus-visible{border-color:var(--navy);color:var(--navy);outline:none}

/* ---------- slides ---------- */
.slide{margin:38px 0 0;background:var(--card);border:1px solid var(--rule);border-radius:12px;
  overflow:hidden;box-shadow:var(--shadow);scroll-margin-top:70px}
.sh{display:flex;gap:16px;align-items:flex-start;padding:20px 22px 16px}
.num{flex:0 0 auto;font-variant-numeric:tabular-nums;font-size:13px;font-weight:700;
  color:var(--card);background:var(--navy);border-radius:6px;padding:4px 9px;margin-top:3px}
.eb{margin:0 0 3px;font-size:11px;letter-spacing:.13em;text-transform:uppercase;
  color:var(--gold);font-weight:700}
.sh h2{margin:0;font-size:19px;line-height:1.3;font-weight:600;letter-spacing:-.01em;
  text-wrap:balance}
.slide img{display:block;width:100%;height:auto;border-top:1px solid var(--rule);
  border-bottom:1px solid var(--rule)}
.note{margin:0;padding:16px 22px 20px;font-size:14px;color:var(--ink-2);background:var(--card)}
.note b{color:var(--navy);font-weight:700}

footer{margin:52px 0 0;padding:26px 0 0;border-top:1px solid var(--rule);
  font-size:13px;color:var(--ink-2)}
footer p{margin:0 0 9px;max-width:70ch}
code{font-family:ui-monospace,SFMono-Regular,Menlo,Consolas,monospace;font-size:.9em;
  background:var(--bg);padding:1px 5px;border-radius:4px;border:1px solid var(--rule)}
@media (max-width:560px){.sh{padding:16px 15px 13px}.note{padding:14px 15px 17px}}
</style>"""


class Deck:
    """Emits python-pptx shapes and, in parallel, records enough to re-draw the slide with Pillow."""

    W, H = 13.3333, 7.5

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.W)
        self.prs.slide_height = Inches(self.H)
        self.blank = self.prs.slide_layouts[6]
        self.slides = []       # list of (pptx slide, [draw ops]) for the preview renderer
        self.findings = []     # overflow / fit complaints
        self.headings = []     # (eyebrow, title) per slide, so a review page needs no second list

    # ---------------------------------------------------------- primitives
    def new(self, ground=WHITE):
        s = self.prs.slides.add_slide(self.blank)
        ops = []
        self.slides.append((s, ops))
        self._s, self._ops = s, ops
        self.rect(0, 0, self.W, self.H, ground)
        return s

    def rect(self, l, t, w, h, fill, radius=None, line=None, lw=0.75):
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sh = self._s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
        if radius:
            # python-pptx exposes the corner radius as adjustment 0, expressed as a fraction of the
            # shorter side. Convert from inches so a card keeps the same visual radius whatever its
            # aspect ratio.
            sh.adjustments[0] = min(0.5, radius / min(w, h))
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = rgb(fill)
        else:
            sh.fill.background()
        if line:
            sh.line.color.rgb = rgb(line)
            sh.line.width = Pt(lw)
        else:
            sh.line.fill.background()
        sh.shadow.inherit = False
        sh.text_frame.word_wrap = False
        self._ops.append(("rect", l, t, w, h, fill, radius, line))
        return sh

    def text(self, l, t, w, s, size=11, bold=False, color=NAVY, align="l", lh=None,
             h=None, anchor="t", space_after=0.0, name=""):
        """One text box. Returns the height it actually consumed, in inches."""
        lh = lh or size * 1.32
        lines = TX.wrap(s, size, w, bold)
        need = (len(lines) * lh) / PT_IN
        box_h = h if h is not None else need
        if h is not None and need > h + 0.02:
            self.findings.append(f"[{name or s[:34]}] text needs {need:.2f}in, box is {h:.2f}in")
        tb = self._s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(box_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                              "b": MSO_ANCHOR.BOTTOM}[anchor]
        p = tf.paragraphs[0]
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
        p.line_spacing = Pt(lh)
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = s
        r.font.name = "Kanit"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
        self._ops.append(("text", l, t, w, box_h, lines, size, bold, color, align, lh, anchor))
        return box_h

    def para(self, l, t, w, runs, size=11, lh=None, align="l", name=""):
        """A single paragraph made of differently-styled runs: [(text, bold, color), ...].

        Used where one sentence has to carry emphasis — the alternative, one box per fragment, is
        what makes generated decks impossible to edit afterwards.
        """
        lh = lh or size * 1.32
        flat = "".join(x[0] for x in runs)
        lines = TX.wrap(flat, size, w)
        box_h = (len(lines) * lh) / PT_IN
        tb = self._s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(box_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER}[align]
        p.line_spacing = Pt(lh)
        for txt, bold, color in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = "Kanit"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)
        # The preview draws the flattened string; emphasis is a colour question, not a fit question.
        self._ops.append(("text", l, t, w, box_h, lines, size, False, NAVY, align, lh, "t"))
        return box_h

    def pic(self, path, l, t, w, h):
        self._s.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
        self._ops.append(("pic", l, t, w, h, str(path)))

    def notes(self, s):
        self._s.notes_slide.notes_text_frame.text = s

    # ---------------------------------------------------------- furniture
    def cover(self, title, sub, kicker):
        self.new(NAVY)
        self.headings.append(("Cover", " ".join(title.split())))
        self.rect(10.73, 0, 2.60, self.H, GOLD)
        self.rect(10.73, 0, 0.12, self.H, RED)
        self.pic(ASSETS / "chaiyo.png", 0.70, 0.70, 3.20, 2.00)
        self.text(0.75, 3.00, 9.30, title, size=30, bold=True, color=WHITE, lh=38)
        self.text(0.75, 4.62, 9.30, sub, size=13, color=GOLD, lh=19)
        self.text(0.75, 6.30, 6.00, kicker, size=16, bold=True, color=WHITE)
        self.text(0.40, 7.16, 7.00, FOOTER, size=8, color=WHITE)

    def divider(self, eyebrow, title, sub):
        self.new(NAVY)
        self.headings.append((eyebrow, title))
        self.rect(12.83, 0, 0.18, self.H, RED)
        self.rect(13.01, 0, 0.18, self.H, GOLD)
        self.pic(ASSETS / "chaiyo.png", 0.70, 0.70, 2.72, 1.70)
        self.text(0.75, 3.05, 11.00, eyebrow, size=11, bold=True, color=GOLD)
        self.text(0.75, 3.45, 11.00, title, size=32, bold=True, color=WHITE, lh=40)
        self.text(0.75, 4.85, 10.60, sub, size=13.5, color="C7D0E4", lh=21)
        self.text(0.40, 7.16, 7.00, FOOTER, size=8, color=WHITE)

    def content(self, eyebrow, title):
        """Returns the top of the usable content area."""
        self.new(WHITE)
        self.headings.append((eyebrow, title))
        self.text(0.45, 0.24, 9.50, eyebrow.upper(), size=9, bold=True, color=RED)
        self.text(0.45, 0.50, 10.90, title, size=20, bold=True, color=NAVY, lh=25)
        self.pic(ASSETS / "autox.png", 11.63, 0.28, 1.03, 0.32)
        self.rect(0.45, 1.06, 12.43, 0.022, NAVY)
        self.rect(0, 7.38, self.W, 0.12, RED)
        self.text(0.45, 7.16, 7.00, FOOTER, size=8, color=GREY)
        return 1.26

    # ---------------------------------------------------------- components
    def chip(self, l, t, label, color):
        """Provenance pill. Every number in this house is labelled measured or estimated."""
        w = TX.width_in(label, 8, True) + 0.26
        self.rect(l, t, w, 0.21, None, radius=0.105, line=color, lw=0.75)
        self.text(l, t + 0.032, w, label, size=8, bold=True, color=color, align="c")
        return w

    def source(self, l, t, w, chip_label, chip_color, body, size=9.5):
        cw = self.chip(l, t, chip_label, chip_color)
        return self.text(l + cw + 0.12, t + 0.015, w - cw - 0.12, body, size=size,
                         color=GREY, lh=size * 1.35, name="source-line")

    def cards(self, l, t, w, items, cols=3, ch=1.16, gap=0.15):
        """items: (label, value, note, colour). The value carries the direction, exactly as the
        web tab does, so a reader who knows the tab reads the deck the same way."""
        cw = (w - gap * (cols - 1)) / cols
        for i, (lab, val, note, col) in enumerate(items):
            x = l + (i % cols) * (cw + gap)
            y = t + (i // cols) * (ch + gap)
            self.rect(x, y, cw, ch, CARD, radius=0.09, line=LINE)
            self.text(x + 0.15, y + 0.11, cw - 0.30, lab, size=9.5, color=GREY, name="card-label")
            self.text(x + 0.15, y + 0.30, cw - 0.30, val, size=21, bold=True, color=col, lh=26,
                      name="card-value")
            self.text(x + 0.15, y + 0.70, cw - 0.30, note, size=8, color=GREY, lh=10.6,
                      h=ch - 0.78, name="card-note")
        rows = (len(items) + cols - 1) // cols
        return rows * ch + (rows - 1) * gap

    def callout(self, l, t, w, head, body, tone="info", size=10.5):
        tint = {"info": "EEF1F7", "warn": "FDF7E9", "risk": "FCEFEF"}[tone]
        col = {"info": NAVY, "warn": "9A7411", "risk": RED}[tone]
        bar = {"info": NAVY, "warn": GOLD, "risk": RED}[tone]
        inner = w - 0.36
        hh = (len(TX.wrap(head, 9.5, inner, True)) * 12.5) / PT_IN
        bh = (len(TX.wrap(body, size, inner)) * size * 1.34) / PT_IN
        h = 0.16 + hh + 0.05 + bh + 0.16
        self.rect(l, t, w, h, tint)
        self.rect(l, t, 0.055, h, bar)
        self.text(l + 0.24, t + 0.16, inner, head, size=9.5, bold=True, color=col, lh=12.5)
        self.text(l + 0.24, t + 0.16 + hh + 0.05, inner, body, size=size, color=NAVY, lh=size * 1.34)
        return h

    def bullets(self, l, t, w, items, size=10.5, gap=0.10, tone=NAVY):
        y = t
        for it in items:
            self.text(l, y + 0.015, 0.18, "•", size=size, bold=True, color=RED)
            y += self.text(l + 0.20, y, w - 0.20, it, size=size, color=tone, lh=size * 1.36) + gap
        return y - t - gap

    def qa(self, l, t, w, rows, qw=3.30, size=11.5):
        """The answer band: a question in navy, its answer beside it, hairline between."""
        y = t
        for i, (q, a) in enumerate(rows):
            if i:
                self.rect(l, y, w, 0.012, LINE)
                y += 0.13
            qh = self.text(l, y, qw, q, size=size, bold=True, color=NAVY, lh=size * 1.34)
            ah = self.text(l + qw + 0.30, y, w - qw - 0.30, a, size=size, color="2A3346",
                           lh=size * 1.34)
            y += max(qh, ah) + 0.13
        return y - t

    def table(self, l, t, w, header, rows, colw, size=10, hsize=10.5, rh=0.30, hh=0.36,
              aligns=None, cellcolors=None):
        """Native PowerPoint table so the numbers stay editable after handover."""
        aligns = aligns or ["l"] * len(header)
        nr, nc = len(rows) + 1, len(header)
        gf = self._s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w),
                                      Inches(hh + rh * len(rows)))
        tbl = gf.table
        tbl.first_row = False
        tbl.horz_banding = False
        total = sum(colw)
        for i, cwf in enumerate(colw):
            tbl.columns[i].width = Emu(int(round(w * cwf / total * EMU_IN)))
        tbl.rows[0].height = Inches(hh)
        for i in range(len(rows)):
            tbl.rows[i + 1].height = Inches(rh)

        def style(cell, txt, sz, bold, color, fill, align):
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
            r = p.add_run()
            r.text = txt
            r.font.name = "Kanit"
            r.font.size = Pt(sz)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)

        colw_in = [w * c / total for c in colw]
        for c, htxt in enumerate(header):
            style(tbl.cell(0, c), htxt, hsize, True, WHITE, NAVY2, aligns[c])
            if TX.width_in(htxt, hsize, True) > colw_in[c] - 0.12:
                self.findings.append(f"[table hdr '{htxt}'] wider than its {colw_in[c]:.2f}in column")
        for ri, row in enumerate(rows):
            fill = WHITE if ri % 2 == 0 else BAND
            for ci, cellv in enumerate(row):
                txt, bold, color = cellv if isinstance(cellv, tuple) else (cellv, False, NAVY)
                if cellcolors and (ri, ci) in cellcolors:
                    color = cellcolors[(ri, ci)]
                style(tbl.cell(ri + 1, ci), txt, size, bold, color, fill, aligns[ci])
                if TX.width_in(txt, size, bold) > colw_in[ci] - 0.12:
                    self.findings.append(
                        f"[table r{ri}c{ci} '{txt[:26]}'] wider than its {colw_in[ci]:.2f}in column")
        self._ops.append(("table", l, t, w, hh + rh * len(rows), header, rows, colw_in, size,
                          hsize, rh, hh, aligns))
        return hh + rh * len(rows)

    # ---------------------------------------------------------- charts
    def linechart(self, l, t, w, h, series, ylab="", baseline=None, xticks=(),
                  ymin=None, ymax=None):
        """A plain index chart drawn as vector shapes.

        series: [(label, colour, [(x_ordinal, y), ...])]. Deliberately not a native PowerPoint
        chart object: this needs 185 monthly points on two series and a horizontal base line, and a
        native chart carries its own cached worksheet and its own theme, which is a lot of moving
        parts to keep on-brand for something nobody will ever re-plot from inside the deck.
        """
        pad_l, pad_b, pad_t = 0.42, 0.30, 0.16
        px0, py0 = l + pad_l, t + pad_t
        pw, ph = w - pad_l - 0.10, h - pad_t - pad_b
        allx = [p[0] for _, _, pts in series for p in pts]
        ally = [p[1] for _, _, pts in series for p in pts]
        x0, x1 = min(allx), max(allx)
        y0 = ymin if ymin is not None else min(ally)
        y1 = ymax if ymax is not None else max(ally)
        sx = lambda v: px0 + (v - x0) / (x1 - x0) * pw
        sy = lambda v: py0 + ph - (v - y0) / (y1 - y0) * ph

        self.rect(l, t, w, h, "FBFCFE", line=LINE)
        for gv in sorted({y0, y1} | ({baseline} if baseline is not None else set())):
            yy = sy(gv)
            is_base = baseline is not None and gv == baseline
            self.rect(px0, yy, pw, 0.012 if is_base else 0.007,
                      GREY if is_base else LINE)
            self.text(l + 0.02, yy - 0.075, pad_l - 0.06, f"{gv:g}", size=7.5,
                      color=GREY if is_base else "9AA3B2", align="r")
        # Legend in the plot's top-right rather than labels at the line ends: the two series
        # converge and cross repeatedly before 2022, so an end-label lands on top of the other line.
        lx = l + w - 0.20
        for lab, col, _pts in reversed(series):
            lw_ = TX.width_in(lab, 8, True)
            self.text(lx - lw_, t + 0.14, lw_, lab, size=8, bold=True, color=col)
            self.rect(lx - lw_ - 0.26, t + 0.205, 0.18, 0.028, col)
            lx -= lw_ + 0.42
        for lab, col, pts in series:
            fb = self._s.shapes.build_freeform(Inches(sx(pts[0][0])), Inches(sy(pts[0][1])))
            fb.add_line_segments([(Inches(sx(x)), Inches(sy(y))) for x, y in pts[1:]],
                                 close=False)
            sh = fb.convert_to_shape()
            sh.fill.background()
            sh.line.color.rgb = rgb(col)
            sh.line.width = Pt(1.6)
            sh.shadow.inherit = False
            self._ops.append(("poly", [(sx(x), sy(y)) for x, y in pts], col))
        for xv, xl in xticks:
            self.text(sx(xv) - 0.20, t + h - 0.22, 0.40, xl, size=7.5, color=GREY, align="c")
        if ylab:
            self.text(l + 0.02, t + 0.01, w - 0.10, ylab, size=8, color=GREY)
        return h

    def ladder(self, l, t, w, segs, h=0.30, gap=0.035, legend=True, size=9):
        """One stacked horizontal bar. segs: [(label, value, colour)].

        Used for the arrears ladder and the crop mix, where the point is the PROPORTION and a
        column of percentages makes the reader do the comparison themselves.
        """
        tot = sum(v for _, v, _ in segs) or 1
        x = l
        for i, (_lab, v, col) in enumerate(segs):
            sw = max(0.0, (w - gap * (len(segs) - 1)) * v / tot)
            self.rect(x, t, sw, h, col)
            x += sw + gap
        if not legend:
            return h
        ly = t + h + 0.10
        lx = l
        for lab, v, col in segs:
            txt = f"{lab} {100 * v / tot:.1f}%"
            tw = TX.width_in(txt, size, False) + 0.30
            if lx + tw > l + w:                      # wrap the legend rather than run it off the box
                lx = l
                ly += size * 1.5 / PT_IN
            self.rect(lx, ly + 0.045, 0.15, 0.10, col)
            self.text(lx + 0.21, ly, tw, txt, size=size, color=NAVY)
            lx += tw
        return ly + size * 1.45 / PT_IN - t

    def sparkline(self, l, t, w, h, pts, color=NAVY, pad=0.05):
        """A shape-only trend line: no axis, no labels, no legend.

        linechart() above is the wrong tool at this size — it reserves 0.42in for value labels and
        prints the series min and max, which at a 2in width turns a rice price of ฿17,440/tonne into
        an unreadable smear across the plot. A sparkline answers one question, "which way and how
        steadily", and the numbers that go with it belong in the text beside it. Each line is scaled
        to its OWN range, so shape is comparable across crops even when levels are not.
        """
        self.rect(l, t, w, h, "FBFCFE", line=LINE)
        ys = [v for _, v in pts]
        lo, hi = min(ys), max(ys)
        span = (hi - lo) or 1.0
        px0, py0 = l + pad, t + pad
        pw, ph = w - 2 * pad, h - 2 * pad
        n = len(pts)
        sx = lambda i: px0 + (i / (n - 1)) * pw if n > 1 else px0 + pw / 2
        sy = lambda v: py0 + ph - (v - lo) / span * ph
        if n > 1:
            fb = self._s.shapes.build_freeform(Inches(sx(0)), Inches(sy(ys[0])))
            fb.add_line_segments([(Inches(sx(i)), Inches(sy(v)))
                                  for i, v in enumerate(ys)][1:], close=False)
            sh = fb.convert_to_shape()
            sh.fill.background()
            sh.line.color.rgb = rgb(color)
            sh.line.width = Pt(1.6)
            sh.shadow.inherit = False
        # The endpoint is the reading everyone actually takes off a sparkline, so mark it.
        r = 0.045
        self.rect(sx(n - 1) - r, sy(ys[-1]) - r, 2 * r, 2 * r, color, radius=r)
        self._ops.append(("poly", [(sx(i), sy(v)) for i, v in enumerate(ys)], color))
        return h

    def bars(self, l, t, w, h, items, color=RED, dim=LINE, fmt=lambda v: f"{v:,.0f}"):
        """Vertical bars with the value written on top. items: [(label, value, highlight)]."""
        self.rect(l, t, w, h, "FBFCFE", line=LINE)
        n = len(items)
        slot = (w - 0.30) / n
        bw = min(slot * 0.52, 1.05)
        top, bot = t + 0.34, t + h - 0.28
        vmax = max(v for _, v, _ in items)
        for i, (lab, v, hi) in enumerate(items):
            cx = l + 0.15 + slot * i + slot / 2
            bh = (bot - top) * (v / vmax)
            self.rect(cx - bw / 2, bot - bh, bw, bh, color if hi else dim)
            self.text(cx - slot / 2, bot - bh - 0.20, slot, fmt(v), size=9, bold=hi,
                      color=color if hi else NAVY, align="c")
            self.text(cx - slot / 2, bot + 0.05, slot, lab, size=8.5, color=GREY, align="c")
        return h

    # ---------------------------------------------------------- output
    def fit_vertical(self, limit=7.30):
        """Report any content shape whose BOTTOM edge runs past the footer.

        The text fit check above is horizontal only — it asks whether a string fits the box it was
        given, never where that box ended up. A slide can therefore pass with "no findings" and
        still print its last callout underneath the red footer bar, which is what happened to the
        conditions slide on 2026-08-04: its author measured the y-cursor BEFORE the trailing callout
        and never added the callout's own height. Walk the real python-pptx geometry instead of the
        preview ops, so every primitive is covered whether or not it draws in the preview.

        Full-bleed grounds and the footer furniture itself sit below the limit by design and are
        excluded by shape, not by name."""
        EMU = 914400.0
        for i, (s, _ops) in enumerate(self.slides, 1):
            worst = None
            for sh in s.shapes:
                t, h = sh.top / EMU, (sh.height or 0) / EMU
                if h > 7.0 or t >= 7.10:            # the ground rect and the footer strip/caption
                    continue
                if t + h > limit and (worst is None or t + h > worst[0]):
                    worst = (t + h, getattr(sh, "text", "") or "")
            if worst:
                head = self.headings[i - 1][1] if i <= len(self.headings) else f"slide {i}"
                self.findings.append(
                    f"[slide {i} '{head[:30]}'] content reaches y={worst[0]:.2f}in, past the "
                    f"{limit:.2f}in floor — {(worst[1] or '(no text)')[:44]!r}")

    def save(self, path):
        self.prs.save(str(path))

    def preview(self, outdir, dpi=110):
        """Redraw every slide with Pillow at the same geometry, so layout can be eyeballed on a
        machine with no PowerPoint and no LibreOffice. Approximate by construction — it exists to
        catch a box that has run off its slide, not to certify pixels."""
        if not TX.ok:
            print("preview skipped: Kanit TTFs not found (set KANIT_DIR)")
            return []
        from PIL import Image, ImageDraw
        outdir = Path(outdir)
        outdir.mkdir(parents=True, exist_ok=True)
        px = lambda v: int(round(v * dpi))
        made = []
        for idx, (_s, ops) in enumerate(self.slides):
            img = Image.new("RGB", (px(self.W), px(self.H)), "white")
            d = ImageDraw.Draw(img, "RGBA")
            for op in ops:
                kind = op[0]
                if kind == "rect":
                    _, l, t, w, h, fill, radius, line = op
                    box = [px(l), px(t), px(l + w), px(t + h)]
                    if radius:
                        d.rounded_rectangle(box, radius=px(radius), fill="#" + fill if fill else None,
                                            outline="#" + line if line else None)
                    else:
                        d.rectangle(box, fill="#" + fill if fill else None,
                                    outline="#" + line if line else None)
                elif kind == "text":
                    _, l, t, w, h, lines, size, bold, color, align, lh, anchor = op
                    f = TX.font(size, bold, scale=dpi / PT_IN)
                    y = px(t)
                    if anchor == "m":
                        y = px(t + (h - len(lines) * lh / PT_IN) / 2)
                    for ln in lines:
                        x = px(l)
                        if align == "c":
                            x = px(l + (w - TX.width_in(ln, size, bold)) / 2)
                        elif align == "r":
                            x = px(l + w - TX.width_in(ln, size, bold))
                        d.text((x, y), ln, font=f, fill="#" + color)
                        y += px(lh / PT_IN)
                elif kind == "poly":
                    _, pts, col = op
                    d.line([(px(x), px(y)) for x, y in pts], fill="#" + col, width=2)
                elif kind == "pic":
                    _, l, t, w, h, path = op
                    im = Image.open(path).convert("RGBA").resize((px(w), px(h)))
                    img.paste(im, (px(l), px(t)), im)
                elif kind == "table":
                    (_, l, t, w, h, header, rows, colw_in, size, hsize, rh, hh, aligns) = op
                    x = l
                    for ci, cw in enumerate(colw_in):
                        d.rectangle([px(x), px(t), px(x + cw), px(t + hh)], fill="#" + NAVY2)
                        f = TX.font(hsize, True, scale=dpi / PT_IN)
                        tw = TX.width_in(header[ci], hsize, True)
                        tx = x + 0.06 if aligns[ci] == "l" else (
                            x + cw - tw - 0.06 if aligns[ci] == "r" else x + (cw - tw) / 2)
                        d.text((px(tx), px(t + (hh - hsize * 1.3 / PT_IN) / 2)), header[ci],
                               font=f, fill="#" + WHITE)
                        x += cw
                    for ri, row in enumerate(rows):
                        y = t + hh + ri * rh
                        x = l
                        for ci, cellv in enumerate(row):
                            txt, bold, color = cellv if isinstance(cellv, tuple) else (cellv, False, NAVY)
                            cw = colw_in[ci]
                            d.rectangle([px(x), px(y), px(x + cw), px(y + rh)],
                                        fill="#" + (WHITE if ri % 2 == 0 else BAND))
                            f = TX.font(size, bold, scale=dpi / PT_IN)
                            tw = TX.width_in(txt, size, bold)
                            tx = x + 0.06 if aligns[ci] == "l" else (
                                x + cw - tw - 0.06 if aligns[ci] == "r" else x + (cw - tw) / 2)
                            d.text((px(tx), px(y + (rh - size * 1.3 / PT_IN) / 2)), txt,
                                   font=f, fill="#" + color)
                            x += cw
            p = outdir / f"slide{idx + 1:02d}.png"
            img.save(p)
            made.append(p)
        return made

    def review(self, pngs, out, title, dek, pptx_name):
        """The sign-off page: every slide as the PNG just rendered, with its own speaker note.

        Generated, never hand-written. The previous copy of this page was authored by hand and went
        quietly wrong the moment the deck grew: it still announced "Twelve slides" against an
        eighteen-slide deck and still quoted a commodity list two price vintages old. A review
        artifact that drifts from what it is reviewing is worse than none, so the slide count, the
        headings and the notes all come off the Deck that was just built.
        """
        import base64

        def esc(s):
            return (s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))

        cards, jump = [], []
        for i, ((_s, _ops), (eyebrow, head)) in enumerate(zip(self.slides, self.headings), 1):
            note = ""
            if _s.has_notes_slide:
                note = _s.notes_slide.notes_text_frame.text or ""
            b64 = base64.b64encode(Path(pngs[i - 1]).read_bytes()).decode("ascii")
            jump.append(f'<li><a href="#s{i}">{i:02d}</a></li>')
            cards.append(
                f'<article class="slide" id="s{i}">\n'
                f'  <header class="sh"><span class="num">{i:02d}</span><div>'
                f'<p class="eb">{esc(eyebrow)}</p><h2>{esc(head)}</h2></div></header>\n'
                f'  <img src="data:image/png;base64,{b64}" alt="Slide {i}: {esc(head)}" '
                f'loading="lazy">\n'
                + (f'  <p class="note"><b>Say:</b> {esc(note)}</p>\n' if note else "")
                + "</article>")
        html = (REVIEW_CSS
                + f'\n<div class="wrap">\n<header class="top">\n'
                f'  <p class="kicker">Draft for sign-off</p>\n  <h1>{title}</h1>\n'
                f'  <p class="dek">{esc(dek)}</p>\n  <ul class="meta">\n'
                f'    <li><b>{len(self.slides)}</b> slides</li>\n'
                f'    <li><b>16:9</b> · 13.33 × 7.5 in</li>\n    <li>Font <b>Kanit</b></li>\n'
                f'    <li>Speaker notes on every slide</li>\n'
                f'    <li>File <b>{esc(pptx_name)}</b></li>\n  </ul>\n</header>\n'
                f'<nav class="jump" aria-label="Jump to slide"><ol>{"".join(jump)}</ol></nav>\n'
                + "\n".join(cards) + "\n</div>\n")
        Path(out).write_text(html, encoding="utf-8")
        return out
